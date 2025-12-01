from pptx import Presentation
from pptx.util import Pt, Inches
from pathlib import Path
md_path = Path('TECHNICAL_REPORT.md')
if not md_path.exists():
    print('TECHNICAL_REPORT.md not found')
    raise SystemExit(1)
md = md_path.read_text(encoding='utf-8')
# Split into sections by the horizontal rule used in the report
sections = []
parts = md.split('\n\n---\n\n')
for part in parts:
    part = part.strip()
    if part:
        sections.append(part)
prs = Presentation()
# Title slide
title_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_layout)
slide.shapes.title.text = 'Gender & Age Detection — Technical Report'
slide.placeholders[1].text = 'Generated from project TECHNICAL_REPORT.md'
# For each section create a slide with bullets
for sec in sections:
    lines = sec.splitlines()
    if not lines:
        continue
    header = lines[0]
    body_lines = [l.strip() for l in lines[1:] if l.strip()]
    # choose a simple layout
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = header[:60]
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for bl in body_lines:
        # truncate long lines
        text = bl if len(bl) <= 240 else bl[:237] + '...'
        p = tf.add_paragraph()
        p.text = text
        p.level = 0
# save presentation
out = Path('TECHNICAL_REPORT.pptx')
prs.save(str(out))
print(f'WROTE: {out.resolve()}')
